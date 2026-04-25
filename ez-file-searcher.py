import os
import threading
import subprocess
import sys
import time
import textwrap
import customtkinter as ctk
from tkinter import messagebox

# File processing libraries
import docx
import openpyxl
import fitz  # PyMuPDF
fitz.TOOLS.mupdf_display_errors(False) 
import pptx  # python-pptx

# --- Setup CustomTkinter Theme ---
ctk.set_appearance_mode("Dark")
ctk.set_default_color_theme("blue")

class EZFileSearcherApp(ctk.CTk):
    def __init__(self):
        super().__init__()

        self.title("EZ-FileSearcher v1.1")
        self.geometry("1100x650")
        self.minsize(800, 500)

        # State variables
        self.selected_folder = ctk.StringVar()
        self.search_term = ctk.StringVar()
        self.is_searching = False
        self.results_data = {} 
        self.current_preview_file = None
        
        # Timing and Navigation state
        self.search_start_time = 0
        self.current_match_indices = [] 
        self.current_match_position = 0 
        
        self.cancel_search = False # <-- ADD THIS CANCEL FLAG

        self.setup_ui()

    def setup_ui(self):
        # --- Top Bar (Folder + Search) ---
        self.top_frame = ctk.CTkFrame(self)
        self.top_frame.pack(side="top", fill="x", padx=10, pady=10)

        self.btn_browse = ctk.CTkButton(self.top_frame, text="Select Folder", command=self.browse_folder, width=120)
        self.btn_browse.pack(side="left", padx=(0, 10))

        # 1. Update the Folder Box (Removed textvariable & readonly)
        self.entry_path = ctk.CTkEntry(self.top_frame, 
                                       placeholder_text="Target directory path...", 
                                       width=300)
        self.entry_path.pack(side="left", fill="x", expand=True, padx=(0, 10))

        # 2. Update the Search Box (Removed textvariable)
        self.entry_search = ctk.CTkEntry(self.top_frame, 
                                         placeholder_text="Search document contents...", 
                                         width=250)
        self.entry_search.pack(side="left", padx=(0, 10))
        self.entry_search.bind("<Return>", lambda event: self.start_search())

        self.btn_search = ctk.CTkButton(self.top_frame, text="Search", command=self.start_search, width=100)
        self.btn_search.pack(side="left", padx=(0, 10)) # <-- Added 10px padding on the right

        # --- Main Body (Left Panel + Right Panel) ---
        self.main_frame = ctk.CTkFrame(self)
        self.main_frame.pack(side="top", fill="both", expand=True, padx=10, pady=(0, 10))

        # Left Panel (Results)
        self.left_panel = ctk.CTkScrollableFrame(self.main_frame, width=300, label_text="Results")
        self.left_panel.pack(side="left", fill="y", padx=(0, 10))

        # Right Panel (Preview)
        self.right_panel = ctk.CTkFrame(self.main_frame)
        self.right_panel.pack(side="left", fill="both", expand=True)

        # Preview Header for Nav Buttons
        self.preview_header = ctk.CTkFrame(self.right_panel, fg_color="transparent")
        self.preview_header.pack(side="top", fill="x", padx=10, pady=(10, 0))

        self.preview_label = ctk.CTkLabel(self.preview_header, text="File Preview", font=("Arial", 14, "bold"))
        self.preview_label.pack(side="left")

        self.btn_down = ctk.CTkButton(self.preview_header, text="▼", width=30, command=self.next_match, state="disabled")
        self.btn_down.pack(side="right", padx=(5, 0))
        
        self.btn_up = ctk.CTkButton(self.preview_header, text="▲", width=30, command=self.prev_match, state="disabled")
        self.btn_up.pack(side="right")

        self.match_nav_label = ctk.CTkLabel(self.preview_header, text="")
        self.match_nav_label.pack(side="right", padx=(0, 10))

        self.textbox_preview = ctk.CTkTextbox(self.right_panel, wrap="word")
        self.textbox_preview.pack(side="top", fill="both", expand=True, padx=10, pady=10)
        
        # Configure highlighting tags
        self.textbox_preview.tag_config("highlight", background="yellow", foreground="black")
        self.textbox_preview.tag_config("active_highlight", background="orange", foreground="black")

        self.btn_open_file = ctk.CTkButton(self.right_panel, text="Open File", command=self.open_current_file, state="disabled")
        self.btn_open_file.pack(side="bottom", anchor="e", padx=10, pady=(0, 10))

       # --- NEW: Bottom Status Bar w/ Progress ---
        self.status_container = ctk.CTkFrame(self, fg_color="transparent")
        self.status_container.pack(side="bottom", fill="x", padx=10, pady=(0, 5))

        self.status_var = ctk.StringVar(value="Ready")
        self.status_label = ctk.CTkLabel(self.status_container, textvariable=self.status_var, anchor="w", text_color="gray")
        self.status_label.pack(side="left")

        # 1. Add the Cancel button here (Notice we do NOT use .pack() yet!)
        self.btn_cancel = ctk.CTkButton(self.status_container, text="Cancel", command=self.trigger_cancel, 
                                        width=70, height=24, fg_color="#b30000", hover_color="#800000")

        self.progress_bar = ctk.CTkProgressBar(self.status_container, width=200)
        self.progress_bar.set(0) # 0 to 1 scale
        self.progress_bar.pack_forget() # Hide it initially

    def browse_folder(self):
        folder = ctk.filedialog.askdirectory(title="Select Folder to Search")
        if folder:
            self.entry_path.delete(0, "end") # Clear any existing text
            self.entry_path.insert(0, folder) # Insert the new folder path

    def start_search(self):
        folder = self.entry_path.get().strip()   # Grab text straight from UI
        term = self.entry_search.get().strip()   # Grab text straight from UI

        if not folder or not term:
            messagebox.showwarning("Missing Information", "Please select a folder and enter a search term.")
            return
        if self.is_searching:
            return

        self.is_searching = True
        self.search_start_time = time.time()
        self.btn_search.configure(state="disabled", text="Searching...")
        self.status_var.set("Pre-scanning folder...")
        
        # Reset UI
        self.cancel_search = False # Make sure this is reset!
        self.results_data.clear()
        self.current_preview_file = None
        self.textbox_preview.delete("1.0", "end")
        self.btn_open_file.configure(state="disabled")
        self.match_nav_label.configure(text="")
        self.btn_up.configure(state="disabled")
        self.btn_down.configure(state="disabled")
        
        # 2. Show the progress bar AND the Cancel button
        self.btn_cancel.configure(text="Cancel", state="normal")
        self.btn_cancel.pack(side="right", padx=(10, 0)) # Pack cancel button on far right
        self.progress_bar.set(0)
        self.progress_bar.pack(side="right", pady=5) # Pack progress bar beside it

        for widget in self.left_panel.winfo_children():
            widget.destroy()

        thread = threading.Thread(target=self.run_search_thread, args=(folder, term), daemon=True)
        thread.start()

    # UPDATE inside start_search(): Right after self.is_searching = True, add these two lines:
    # self.cancel_search = False
    # self.btn_cancel.configure(state="normal", text="Cancel")

    def trigger_cancel(self):
        if self.is_searching:
            self.cancel_search = True
            self.btn_cancel.configure(state="disabled", text="Stopping...")
            self.status_var.set("Canceling search... (finishing current file)")

    def read_file_content(self, filepath):
        ext = os.path.splitext(filepath)[1].lower()
        plain_text_exts = ['.txt', '.csv', '.py', '.log', '.json', '.md', '.xml', '.html', '.ini', '.cfg', '.bat', '.sh']
        
        try:
            if ext in plain_text_exts:
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f: return f.read()
            elif ext == '.docx':
                return '\n'.join([p.text for p in docx.Document(filepath).paragraphs])
            elif ext == '.xlsx':
                lines = []
                for sheet in openpyxl.load_workbook(filepath, data_only=True).worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        line = ' '.join([str(cell) for cell in row if cell is not None])
                        if line.strip(): lines.append(line)
                return '\n'.join(lines)
            elif ext == '.pdf':
                text = []
                with fitz.open(filepath) as doc:
                    for page in doc: text.append(page.get_text())
                return '\n'.join(text)
            elif ext == '.pptx':
                text = []
                for slide in pptx.Presentation(filepath).slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"): text.append(shape.text)
                return '\n'.join(text)
        except Exception:
            pass
        return None

    def run_search_thread(self, folder, term):
        exts = ('.txt', '.csv', '.py', '.log', '.json', '.md', '.xml', '.html', '.ini', '.cfg', '.bat', '.sh', '.docx', '.xlsx', '.pdf', '.pptx')
        term_lower = term.lower()
        
        # PHASE 1: Lightning Pre-Scan to count files
        target_files = []
        for root, dirs, files in os.walk(folder):
            if self.cancel_search: break # <-- ADD THIS LINE
            for file in files:
                if file.lower().endswith(exts):
                    target_files.append(os.path.join(root, file))
                    
        total_files = len(target_files)
        if total_files == 0:
            self.after(0, self.finish_search, 0, 0, term)
            return

        # PHASE 2: Deep Scan with Progress Tracking
        total_matches = 0
        for index, filepath in enumerate(target_files):
            if self.cancel_search: break # <-- ADD THIS LINE
            content = self.read_file_content(filepath)
            if content:
                match_count = content.lower().count(term_lower)
                if match_count > 0:
                    self.results_data[filepath] = {"count": match_count, "content": content}
                    total_matches += match_count
                    self.after(0, self.add_result_button, filepath, match_count)
            
            # Math for progress bar
            files_processed = index + 1
            progress_float = files_processed / total_files
            percent = int(progress_float * 100)
            
            # Safely update GUI
            self.after(0, self.update_progress_ui, progress_float, percent, files_processed, total_files)

        self.after(0, self.finish_search, total_files, total_matches, term)

    def update_progress_ui(self, progress_float, percent, current, total):
        """Updates the progress bar and the status text."""
        self.progress_bar.set(progress_float)
        self.status_var.set(f"Searching: {percent}% ({current}/{total} files processed)")

    def add_result_button(self, filepath, count):
        filename = os.path.basename(filepath)
        display_text = f"{filename} ({count} matches)"
        
        # Wrap the text at 33 characters (fits nicely in 300px panel)
        wrapped_text = textwrap.fill(display_text, width=33)
        
        # Create a container frame that acts as our "button" background
        btn_frame = ctk.CTkFrame(self.left_panel, fg_color="transparent", cursor="hand2")
        btn_frame.pack(side="top", fill="x", pady=2)
        
        # Create a Label inside the frame (Labels DO support justify="left")
        lbl = ctk.CTkLabel(btn_frame, text=wrapped_text, justify="left", anchor="w", 
                           text_color=("gray10", "gray90"))
        lbl.pack(side="left", fill="x", expand=True, padx=8, pady=4)
        
        # Create manual hover and click events
        def on_enter(e): btn_frame.configure(fg_color=("gray70", "gray30"))
        def on_leave(e): btn_frame.configure(fg_color="transparent")
        def on_click(e, p=filepath): self.show_preview(p)
        
        # Bind the events to both the frame AND the label so it feels like one solid button
        btn_frame.bind("<Enter>", on_enter)
        btn_frame.bind("<Leave>", on_leave)
        btn_frame.bind("<Button-1>", on_click)
        lbl.bind("<Enter>", on_enter)
        lbl.bind("<Leave>", on_leave)
        lbl.bind("<Button-1>", on_click)

    def finish_search(self, files_scanned, total_matches, term):
        self.is_searching = False
        self.btn_search.configure(state="normal", text="Search")
        
        # 3. Hide both widgets when finished
        self.progress_bar.pack_forget() 
        self.btn_cancel.pack_forget()   
        
        elapsed_time = time.time() - self.search_start_time 
        
        if self.cancel_search:
            self.status_var.set(f"Search CANCELED. Found {total_matches} matches so far. ({elapsed_time:.2f}s)")
            return

        if not self.results_data:
            self.status_var.set(f"No files found for '{term}' (Scanned {files_scanned} files in {elapsed_time:.2f}s)")
            lbl = ctk.CTkLabel(self.left_panel, text="No results found.", text_color="gray")
            lbl.pack(pady=20)
        else:
            self.status_var.set(f"Found {total_matches} matches across {len(self.results_data)} files. (Scanned {files_scanned} files in {elapsed_time:.2f}s)")

    def show_preview(self, filepath):
        self.current_preview_file = filepath
        self.btn_open_file.configure(state="normal")
        self.preview_label.configure(text=f"Preview: {os.path.basename(filepath)}")
        
        content = self.results_data[filepath]["content"]
        term = self.entry_search.get().strip()

        self.textbox_preview.delete("1.0", "end")
        self.textbox_preview.insert("1.0", content)

        self.current_match_indices.clear()
        self.current_match_position = 0

        start_idx = "1.0"
        term_len = len(term)
        while True:
            pos = self.textbox_preview.search(term, start_idx, nocase=True, stopindex="end")
            if not pos:
                break
            
            end_idx = f"{pos}+{term_len}c"
            self.textbox_preview.tag_add("highlight", pos, end_idx)
            self.current_match_indices.append((pos, end_idx)) 
            start_idx = end_idx

        if self.current_match_indices:
            self.btn_up.configure(state="normal")
            self.btn_down.configure(state="normal")
            self.update_match_navigation()

    def update_match_navigation(self):
        if not self.current_match_indices: return
        self.textbox_preview.tag_remove("active_highlight", "1.0", "end")
        pos, end_idx = self.current_match_indices[self.current_match_position]
        self.textbox_preview.tag_add("active_highlight", pos, end_idx)
        self.textbox_preview.see(pos)
        
        total = len(self.current_match_indices)
        current = self.current_match_position + 1
        self.match_nav_label.configure(text=f"{current} / {total}")

    def next_match(self):
        if self.current_match_indices:
            self.current_match_position = (self.current_match_position + 1) % len(self.current_match_indices)
            self.update_match_navigation()

    def prev_match(self):
        if self.current_match_indices:
            self.current_match_position = (self.current_match_position - 1) % len(self.current_match_indices)
            self.update_match_navigation()

    def open_current_file(self):
        if not self.current_preview_file: return
        filepath = self.current_preview_file
        if sys.platform == "win32":
            try: os.startfile(filepath)
            except Exception as e: messagebox.showerror("Error", f"Could not open:\n{str(e)}")
        elif sys.platform == "darwin":
            subprocess.call(('open', filepath))
        else:
            subprocess.call(('xdg-open', filepath))

if __name__ == "__main__":
    app = EZFileSearcherApp()
    app.mainloop()
